#!/usr/bin/env python3
# Regenerate frames/ from the source AllBirds PowerPoint deck (one slide per
# species, 3 stacked photo+credit pairs per slide cycled via PER_BIRD).
# Requires Windows + PowerPoint installed (uses COM automation - python-pptx
# alone can edit shapes but can't rasterize slides).
#
# The deck is native 7.5x10in (3:4). This resizes every slide to 9:16
# (matching the Acer VG240Y's 1080x1920 portrait viewport), growing the
# photo panel to absorb the added height and shifting everything below it
# down to match -- so the export fills the screen with no crop or pad,
# and no text box changes size or wraps differently.
#
# Usage: py tools\export_frames_from_pptx.py path\to\AllBirds.pptx
import sys, os
from pptx import Presentation
import win32com.client
from PIL import Image

TARGET_W, TARGET_H = 1080, 1920

def resize_deck(src_path, out_pptx_path):
    prs = Presentation(src_path)
    orig_h = prs.slide_height
    new_h = round(prs.slide_width * TARGET_H / TARGET_W)
    extra = new_h - orig_h

    photo_top, photo_h = 914400, 4114800   # fixed layout constants in this template
    credit_top = 4736592
    lower_threshold = 5138928

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.top == 0 and shape.height == orig_h:
                shape.height = new_h
            elif shape.top == photo_top and shape.height == photo_h:
                shape.height = photo_h + extra
                if shape.shape_type == 13:  # PICTURE
                    # Original crop % was a cover-fit calibrated for the old,
                    # much less tall box -- recompute it as a cover-fit for
                    # the new box so the photo fills it with no distortion
                    # (crop the overflow, same idea as CSS object-fit: cover,
                    # rather than stretching pixels non-uniformly).
                    img_w, img_h = shape.image.size
                    img_aspect = img_w / img_h
                    box_aspect = shape.width / shape.height
                    if img_aspect > box_aspect:
                        visible_frac = box_aspect / img_aspect
                        side = (1 - visible_frac) / 2
                        shape.crop_left = shape.crop_right = side
                        shape.crop_top = shape.crop_bottom = 0.0
                    else:
                        visible_frac = img_aspect / box_aspect
                        side = (1 - visible_frac) / 2
                        shape.crop_top = shape.crop_bottom = side
                        shape.crop_left = shape.crop_right = 0.0
            elif shape.top == credit_top:
                shape.top = credit_top + extra
            elif shape.top >= lower_threshold:
                shape.top = shape.top + extra

    prs.slide_height = new_h
    prs.save(out_pptx_path)

def export_frames(pptx_path, frames_dir):
    app = win32com.client.DispatchEx("PowerPoint.Application")
    app.Visible = True
    pres = app.Presentations.Open(pptx_path, WithWindow=True)

    frame_idx = 0
    for si in range(1, pres.Slides.Count + 1):
        slide = pres.Slides(si)
        pics, credits = [], []
        for i in range(1, slide.Shapes.Count + 1):
            s = slide.Shapes.Item(i)
            if s.Type == 13:  # msoPicture
                pics.append(s)
            elif s.Name.startswith("TextBox") and 600 < s.Top < 620:
                credits.append(s)  # credit caption sits near the photo bottom

        for variant in range(len(pics)):
            for j in range(len(pics)):
                pics[j].Visible = (j == variant)
            for j in range(len(credits)):
                credits[j].Visible = (j == variant)
            png_path = os.path.join(frames_dir, "frame_%03d.png" % frame_idx)
            slide.Export(png_path, "PNG", TARGET_W, TARGET_H)
            frame_idx += 1

        for s in pics + credits:
            s.Visible = True

    pres.Close()
    app.Quit()
    return frame_idx

def pngs_to_jpgs(frames_dir):
    import glob
    for f in sorted(glob.glob(os.path.join(frames_dir, "*.png"))):
        im = Image.open(f).convert("RGB")
        im.save(os.path.splitext(f)[0] + ".jpg", quality=92)
        os.remove(f)

if __name__ == "__main__":
    src = sys.argv[1]
    here = os.path.dirname(os.path.abspath(__file__))
    frames_dir = os.path.join(here, "..", "frames")
    resized_pptx = os.path.join(here, "_resized_9x16.pptx")

    resize_deck(src, resized_pptx)
    n = export_frames(resized_pptx, frames_dir)
    pngs_to_jpgs(frames_dir)
    os.remove(resized_pptx)
    print(f"exported {n} frames to {frames_dir}")
